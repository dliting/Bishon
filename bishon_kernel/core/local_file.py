import email
import email.policy
import os
from collections.abc import Callable

import cv2
import docx
import numpy as np
import pandas as pd
import pptx
from langchain_community.document_loaders import Docx2txtLoader, TextLoader
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx.enum.shapes import MSO_SHAPE_TYPE

from bishon_kernel.configs.model_config import SENTENCE_SIZE, ZH_TITLE_ENHANCE
from bishon_kernel.utils.custom_log import debug_logger
from bishon_kernel.utils.general_utils import num_tokens, write_check_file
from bishon_kernel.utils.loader import UnstructuredPaddleImageLoader, UnstructuredPaddlePDFLoader
from bishon_kernel.utils.loader.csv_loader import CSVLoader
from bishon_kernel.utils.loader.my_recursive_url_loader import MyRecursiveUrlLoader
from bishon_kernel.utils.ocr_utils import numpy_to_ocr_data
from bishon_kernel.utils.splitter import ChineseTextSplitter, zh_title_enhance

TEXT_SPLITTER_CHUNK_SIZE = 400

text_splitter = RecursiveCharacterTextSplitter(
    separators=["\n", ".", "。", "!", "！", "?", "？", "；", ";", "……", "…", "、", "，", ",", " "],
    chunk_size=TEXT_SPLITTER_CHUNK_SIZE,
    length_function=num_tokens,
)

_SUPPORTED_TYPES = "文件类型不支持，目前仅支持：[md,txt,pdf,jpg,png,jpeg,docx,xlsx,pptx,eml,csv]"
_IMAGE_EXTENSIONS = (".jpg", ".png", ".jpeg")


def _ocr_image_bytes(ocr_engine: Callable, img_bytes: bytes) -> list[str]:
    """Run PaddleOCR on raw image bytes, return extracted text lines.

    Invalid/undecodable images return [] silently; OCR engine errors
    propagate to the caller for logging.
    """
    img_np = cv2.imdecode(np.frombuffer(img_bytes, dtype=np.uint8), cv2.IMREAD_COLOR)
    if img_np is None:
        return []
    img_data = numpy_to_ocr_data(img_np)
    del img_np
    result = ocr_engine(img_data)
    return list(result) if result else []


class LocalFile:
    def __init__(self, user_id, kb_id, file, file_id, file_name, embedding, is_url=False, in_milvus=False):
        self.user_id   = user_id
        self.kb_id     = kb_id
        self.file_id   = file_id
        self.docs: list[Document] = []
        self.embs      = []
        self.emb_infer = embedding
        self.url       = None
        self.in_milvus = in_milvus
        self.file_name = file_name
        if is_url:
            self.url       = file
            self.file_path = "URL"
        elif isinstance(file, str):
            self.file_path = file
        else:
            raise TypeError(f"Expected file_path (str) or URL, got {type(file).__name__}")
        debug_logger.info('success init localfile %s', self.file_name)

    # ------------------------------------------------------------------
    # Public entry: dispatch + shared post-processing
    # ------------------------------------------------------------------

    def split_file_to_docs(self, ocr_engine: Callable | None, sentence_size=SENTENCE_SIZE,
                           using_zh_title_enhance=ZH_TITLE_ENHANCE):
        file_path_lower = self.file_path.lower()

        if self.url:
            docs = self._load_url(sentence_size)
        elif file_path_lower.endswith((".md", ".txt")):
            docs = self._load_text(sentence_size)
        elif file_path_lower.endswith(".pdf"):
            docs = self._load_pdf(ocr_engine, sentence_size)
        elif file_path_lower.endswith(_IMAGE_EXTENSIONS):
            docs = self._load_image(ocr_engine, sentence_size)
        elif file_path_lower.endswith(".docx"):
            docs = self._load_docx(ocr_engine, sentence_size)
        elif file_path_lower.endswith(".xlsx"):
            docs = self._load_xlsx()
        elif file_path_lower.endswith(".pptx"):
            docs = self._load_pptx(ocr_engine)
        elif file_path_lower.endswith(".eml"):
            docs = self._load_eml()
        elif file_path_lower.endswith(".csv"):
            docs = self._load_csv()
        else:
            raise TypeError(_SUPPORTED_TYPES)

        if using_zh_title_enhance:
            debug_logger.info("using_zh_title_enhance %s", using_zh_title_enhance)
            docs = zh_title_enhance(docs)

        debug_logger.info("before 2nd split doc lens: %d", len(docs))
        docs = text_splitter.split_documents(docs)
        debug_logger.info("after 2nd split doc lens: %d", len(docs))

        for doc in docs:
            doc.metadata["file_id"]   = self.file_id
            doc.metadata["file_name"] = self.url if self.url else os.path.split(self.file_path)[-1]
        write_check_file(self.file_path, docs)
        if docs:
            debug_logger.info('langchain analysis content head: %s', docs[0].page_content[:100])
        else:
            debug_logger.info('langchain analysis docs is empty!')
        self.docs = docs

    # ------------------------------------------------------------------
    # Per-file-type loaders
    # ------------------------------------------------------------------

    def _load_url(self, sentence_size) -> list[Document]:
        debug_logger.info("load url: %s", self.url)
        loader       = MyRecursiveUrlLoader(url=self.url)
        textsplitter = ChineseTextSplitter(pdf=False, sentence_size=sentence_size)
        return loader.load_and_split(text_splitter=textsplitter)

    def _load_text(self, sentence_size) -> list[Document]:
        loader         = TextLoader(self.file_path, encoding='utf-8', autodetect_encoding=True)
        texts_splitter = ChineseTextSplitter(pdf=False, sentence_size=sentence_size)
        return loader.load_and_split(texts_splitter)

    def _load_pdf(self, ocr_engine: Callable, sentence_size) -> list[Document]:
        if ocr_engine is None:
            raise ValueError("OCR engine not available, cannot process PDF files. "
                             "Please check PaddleOCR installation.")
        loader         = UnstructuredPaddlePDFLoader(self.file_path, ocr_engine)
        texts_splitter = ChineseTextSplitter(pdf=True, sentence_size=sentence_size)
        return loader.load_and_split(texts_splitter)

    def _load_image(self, ocr_engine: Callable, sentence_size) -> list[Document]:
        if ocr_engine is None:
            raise ValueError("OCR engine not available, cannot process image files. "
                             "Please check PaddleOCR installation.")
        loader         = UnstructuredPaddleImageLoader(self.file_path, ocr_engine, mode="elements")
        texts_splitter = ChineseTextSplitter(pdf=False, sentence_size=sentence_size)
        return loader.load_and_split(text_splitter=texts_splitter)

    def _load_docx(self, ocr_engine: Callable | None, sentence_size) -> list[Document]:
        loader         = Docx2txtLoader(self.file_path)
        texts_splitter = ChineseTextSplitter(pdf=False, sentence_size=sentence_size)
        docs = loader.load_and_split(texts_splitter)
        if ocr_engine is not None:
            docx_doc  = docx.Document(self.file_path)
            ocr_texts = []
            for rel in docx_doc.part.rels.values():
                if "image" in rel.reltype:
                    try:
                        ocr_texts.extend(_ocr_image_bytes(ocr_engine, rel.target_part.blob))
                    except Exception:
                        debug_logger.warning("OCR failed on image in DOCX %s",
                                             self.file_name, exc_info=True)
            if ocr_texts:
                docs.append(Document(page_content='\n'.join(ocr_texts),
                                     metadata={"source": self.file_path}))
        return docs

    def _load_xlsx(self) -> list[Document]:
        docs           = []
        temp_csv_paths = []
        try:
            xlsx = pd.read_excel(self.file_path, engine='openpyxl', sheet_name=None)
            for sheet in xlsx.keys():
                df             = xlsx[sheet]
                df.dropna(how='all', inplace=True)
                csv_file_path  = self.file_path[:-5] + '_' + sheet + '.csv'
                temp_csv_paths.append(csv_file_path)
                df.to_csv(csv_file_path, index=False)
                loader = CSVLoader(csv_file_path, csv_args={"delimiter": ",", "quotechar": '"'},
                                   autodetect_encoding=True)
                docs += loader.load()
        finally:
            for csv_path in temp_csv_paths:
                try:
                    os.remove(csv_path)
                except OSError:
                    debug_logger.warning("Failed to clean up temp CSV: %s", csv_path)
        return docs

    def _load_pptx(self, ocr_engine: Callable | None) -> list[Document]:
        prs   = pptx.Presentation(self.file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
                if ocr_engine is not None and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        ocr_texts = _ocr_image_bytes(ocr_engine, shape.image.blob)
                        texts.extend(ocr_texts)
                    except Exception:
                        debug_logger.warning("OCR failed on image in PPTX %s",
                                             self.file_name, exc_info=True)
        return [Document(page_content='\n\n'.join(texts), metadata={"source": self.file_path})]

    def _load_eml(self) -> list[Document]:
        with open(self.file_path, 'rb') as f:
            msg = email.message_from_binary_file(f, policy=email.policy.default)
        body_parts = []
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == 'text/plain':
                    payload = part.get_payload(decode=True)
                    if payload:
                        charset = part.get_content_charset() or 'utf-8'
                        body_parts.append(payload.decode(charset, errors='replace'))
        else:
            payload = msg.get_payload(decode=True)
            if payload:
                charset = msg.get_content_charset() or 'utf-8'
                body_parts.append(payload.decode(charset, errors='replace'))
        subject = msg.get('Subject', '')
        if subject:
            body_parts.insert(0, f"Subject: {subject}")
        return [Document(page_content='\n\n'.join(body_parts), metadata={"source": self.file_path})]

    def _load_csv(self) -> list[Document]:
        loader = CSVLoader(self.file_path, csv_args={"delimiter": ",", "quotechar": '"'},
                           autodetect_encoding=True, encoding='utf-8')
        return loader.load()

    # ------------------------------------------------------------------

    def create_embedding(self):
        self.embs = self.emb_infer._get_len_safe_embeddings([doc.page_content for doc in self.docs])
