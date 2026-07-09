"""Tests for LocalFile parsing and splitting."""
import io
import os
import textwrap
from unittest.mock import MagicMock

import numpy as np
import pytest


@pytest.fixture
def mock_embedding():
    emb = MagicMock()
    emb._get_len_safe_embeddings.return_value = [[0.1] * 768]
    return emb


class TestLocalFileTxt:
    def test_txt_split(self, tmp_path, mock_embedding):
        from bishon_kernel.core.local_file import LocalFile
        txt_file = tmp_path / "test.txt"
        txt_file.write_text("Hello world. This is a test document with some content.", encoding="utf-8")
        lf = LocalFile("user1", "KB001", str(txt_file), "f1", "test.txt", mock_embedding)
        lf.split_file_to_docs(ocr_engine=None)
        assert len(lf.docs) > 0
        assert lf.docs[0].metadata["file_id"] == "f1"
        assert lf.docs[0].metadata["file_name"] == "test.txt"

    def test_txt_file_path_set(self, tmp_path, mock_embedding):
        from bishon_kernel.core.local_file import LocalFile
        txt_file = tmp_path / "hello.txt"
        content = "人工智能是计算机科学的一个分支。"
        txt_file.write_text(content, encoding="utf-8")
        lf = LocalFile("user1", "KB001", str(txt_file), "f1", "hello.txt", mock_embedding)
        assert lf.file_path == str(txt_file)
        assert os.path.exists(lf.file_path)
        assert not hasattr(lf, 'file_content')


class TestLocalFileCsv:
    def test_csv_split(self, tmp_path, mock_embedding):
        from bishon_kernel.core.local_file import LocalFile
        csv_file = tmp_path / "test.csv"
        csv_file.write_text("name,age\nAlice,30\nBob,25\n", encoding="utf-8")
        lf = LocalFile("user1", "KB001", str(csv_file), "f2", "test.csv", mock_embedding)
        lf.split_file_to_docs(ocr_engine=None)
        assert len(lf.docs) > 0


class TestLocalFileUnsupported:
    def test_unsupported_type_raises(self, tmp_path, mock_embedding):
        from bishon_kernel.core.local_file import LocalFile
        rar_file = tmp_path / "test.rar"
        rar_file.write_bytes(b"fake rar content")
        lf = LocalFile("user1", "KB001", str(rar_file), "f3", "test.rar", mock_embedding)
        with pytest.raises(TypeError, match="文件类型不支持"):
            lf.split_file_to_docs(ocr_engine=None)


class TestLocalFileCreateEmbedding:
    def test_create_embedding(self, tmp_path, mock_embedding):
        from bishon_kernel.connector.database.faiss.faiss_client import Document
        from bishon_kernel.core.local_file import LocalFile
        txt_file = tmp_path / "test.txt"
        txt_file.write_text("Hello world", encoding="utf-8")
        lf = LocalFile("user1", "KB001", str(txt_file), "f1", "test.txt", mock_embedding)
        lf.docs = [Document(page_content="hello")]
        mock_embedding._get_len_safe_embeddings.return_value = [[0.1] * 768]
        lf.create_embedding()
        assert len(lf.embs) == 1
        assert len(lf.embs[0]) == 768


class TestLocalFileOcrUnavailable:
    """When OCR is unavailable, images and PDFs must raise a clear ValueError (not TypeError)."""

    def test_image_raises_when_ocr_none(self, tmp_path, mock_embedding):
        from bishon_kernel.core.local_file import LocalFile
        img_file = tmp_path / "test.png"
        img_file.write_bytes(b"\x89PNG\r\n\x1a\n" + b"\x00" * 100)
        lf = LocalFile("user1", "KB001", str(img_file), "f1", "test.png", mock_embedding)
        with pytest.raises(ValueError, match="OCR"):
            lf.split_file_to_docs(ocr_engine=None)

    def test_jpg_raises_when_ocr_none(self, tmp_path, mock_embedding):
        from bishon_kernel.core.local_file import LocalFile
        jpg_file = tmp_path / "test.jpg"
        jpg_file.write_bytes(b"\xff\xd8\xff\xe0" + b"\x00" * 100)
        lf = LocalFile("user1", "KB001", str(jpg_file), "f2", "test.jpg", mock_embedding)
        with pytest.raises(ValueError, match="OCR"):
            lf.split_file_to_docs(ocr_engine=None)

    def test_pdf_raises_when_ocr_none(self, tmp_path, mock_embedding):
        from bishon_kernel.core.local_file import LocalFile
        pdf_file = tmp_path / "test.pdf"
        pdf_file.write_bytes(b"%PDF-1.4" + b"\x00" * 100)
        lf = LocalFile("user1", "KB001", str(pdf_file), "f3", "test.pdf", mock_embedding)
        with pytest.raises(ValueError, match="OCR"):
            lf.split_file_to_docs(ocr_engine=None)

    def test_jpeg_raises_when_ocr_none(self, tmp_path, mock_embedding):
        from bishon_kernel.core.local_file import LocalFile
        jpeg_file = tmp_path / "test.jpeg"
        jpeg_file.write_bytes(b"\xff\xd8\xff\xe0" + b"\x00" * 100)
        lf = LocalFile("user1", "KB001", str(jpeg_file), "f4", "test.jpeg", mock_embedding)
        with pytest.raises(ValueError, match="OCR"):
            lf.split_file_to_docs(ocr_engine=None)


class TestLocalFileUrl:
    def test_url_mode(self, mock_embedding):
        from bishon_kernel.core.local_file import LocalFile
        lf = LocalFile("user1", "KB001", "https://example.com", "f1", "https://example.com", mock_embedding, is_url=True)
        assert lf.url == "https://example.com"
        assert lf.file_path == "URL"
        assert not hasattr(lf, 'file_content')


class TestLocalFileMd:
    def test_md_split(self, tmp_path, mock_embedding):
        from bishon_kernel.core.local_file import LocalFile
        md_file = tmp_path / "test.md"
        md_file.write_text("# Title\n\nThis is a markdown document with some content.", encoding="utf-8")
        lf = LocalFile("user1", "KB001", str(md_file), "f_md", "test.md", mock_embedding)
        lf.split_file_to_docs(ocr_engine=None)
        assert len(lf.docs) > 0
        assert lf.docs[0].metadata["file_id"] == "f_md"


class TestLocalFileDocx:
    def test_docx_split(self, tmp_path, mock_embedding):
        from docx import Document as DocxDocument
        docx_path = tmp_path / "test.docx"
        doc = DocxDocument()
        doc.add_paragraph("Hello from docx. This is a test document for Bishon.")
        doc.save(str(docx_path))

        from bishon_kernel.core.local_file import LocalFile
        lf = LocalFile("user1", "KB001", str(docx_path), "f_docx", "test.docx", mock_embedding)
        lf.split_file_to_docs(ocr_engine=None)
        assert len(lf.docs) > 0
        assert lf.docs[0].metadata["file_id"] == "f_docx"
        assert "Hello" in lf.docs[0].page_content


class TestLocalFileXlsx:
    def test_xlsx_split(self, tmp_path, mock_embedding):
        from openpyxl import Workbook
        xlsx_path = tmp_path / "test.xlsx"
        wb = Workbook()
        ws = wb.active
        ws.append(["Name", "Age"])
        ws.append(["Alice", "30"])
        ws.append(["Bob", "25"])
        wb.save(str(xlsx_path))

        from bishon_kernel.core.local_file import LocalFile
        lf = LocalFile("user1", "KB001", str(xlsx_path), "f_xlsx", "test.xlsx", mock_embedding)
        lf.split_file_to_docs(ocr_engine=None)
        assert len(lf.docs) > 0
        assert lf.docs[0].metadata["file_id"] == "f_xlsx"


class TestLocalFilePptx:
    def test_pptx_split(self, tmp_path, mock_embedding):
        from pptx import Presentation
        pptx_path = tmp_path / "test.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Test Slide Title"
        slide.placeholders[1].text = "This is the body of a PowerPoint slide."
        prs.save(str(pptx_path))

        from bishon_kernel.core.local_file import LocalFile
        lf = LocalFile("user1", "KB001", str(pptx_path), "f_pptx", "test.pptx", mock_embedding)
        lf.split_file_to_docs(ocr_engine=None)
        assert len(lf.docs) > 0
        assert lf.docs[0].metadata["file_id"] == "f_pptx"


class TestLocalFileEml:
    def test_eml_split(self, tmp_path, mock_embedding):
        eml_content = textwrap.dedent("""\
            From: sender@example.com
            To: receiver@example.com
            Subject: Test Email
            Date: Sat, 17 May 2026 12:00:00 +0800
            MIME-Version: 1.0
            Content-Type: text/plain; charset="utf-8"

            This is the body of a test email for Bishon document parsing.
        """)
        eml_path = tmp_path / "test.eml"
        eml_path.write_text(eml_content, encoding="utf-8")

        from bishon_kernel.core.local_file import LocalFile
        lf = LocalFile("user1", "KB001", str(eml_path), "f_eml", "test.eml", mock_embedding)
        lf.split_file_to_docs(ocr_engine=None)
        assert len(lf.docs) > 0
        assert lf.docs[0].metadata["file_id"] == "f_eml"


class TestLocalFileCsvChinese:
    def test_csv_chinese_split(self, tmp_path, mock_embedding):
        csv_content = "姓名,年龄\n张三,30\n李四,25\n"
        csv_path = tmp_path / "chinese.csv"
        csv_path.write_text(csv_content, encoding="utf-8")

        from bishon_kernel.core.local_file import LocalFile
        lf = LocalFile("user1", "KB001", str(csv_path), "f_csv_cn", "chinese.csv", mock_embedding)
        lf.split_file_to_docs(ocr_engine=None)
        assert len(lf.docs) > 0
        all_content = " ".join(doc.page_content for doc in lf.docs)
        assert "张三" in all_content


def _make_valid_png_bytes():
    """Create a minimal valid 2x2 red PNG image bytes."""
    img = np.zeros((2, 2, 3), dtype=np.uint8)
    img[:, :, 2] = 255
    import cv2
    _, png = cv2.imencode('.png', img)
    return png.tobytes()


class TestPptxImageOcr:
    """PPT image OCR: with ocr_engine, image text is extracted; without, images are skipped."""

    OCR_PREFIX = "[OCR_IMG]"

    def _mock_ocr(self, img_data):
        return [f"{self.OCR_PREFIX} mock text from image"]

    def test_pptx_without_ocr_still_extracts_text(self, tmp_path, mock_embedding):
        """Without OCR, text is still extracted without errors."""
        from pptx import Presentation
        pptx_path = tmp_path / "text_only.pptx"
        prs = Presentation()
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        txBox = slide.shapes.add_textbox(0, 0, 1000000, 1000000)
        txBox.text_frame.text = "Slide text content."
        prs.save(str(pptx_path))

        from bishon_kernel.core.local_file import LocalFile
        lf = LocalFile("user1", "KB001", str(pptx_path), "f1", "text_only.pptx", mock_embedding)
        lf.split_file_to_docs(ocr_engine=None)
        all_text = " ".join(d.page_content for d in lf.docs)
        assert "Slide text content" in all_text

    def test_pptx_with_image_ocr(self, tmp_path, mock_embedding):
        """With OCR, image text is extracted and merged into the result."""
        from pptx import Presentation
        from pptx.util import Inches
        pptx_path = tmp_path / "with_image.pptx"
        prs = Presentation()
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        txBox = slide.shapes.add_textbox(0, 0, 1000000, 1000000)
        txBox.text_frame.text = "Pre-image text."
        img_bytes = _make_valid_png_bytes()
        slide.shapes.add_picture(io.BytesIO(img_bytes), Inches(1), Inches(2), Inches(2), Inches(2))
        txBox2 = slide.shapes.add_textbox(0, Inches(4.5) * 914400, 1000000, 1000000)
        txBox2.text_frame.text = "Post-image text."
        prs.save(str(pptx_path))

        from bishon_kernel.core.local_file import LocalFile
        lf = LocalFile("user1", "KB001", str(pptx_path), "f2", "with_image.pptx", mock_embedding)
        lf.split_file_to_docs(ocr_engine=self._mock_ocr)
        all_text = " ".join(d.page_content for d in lf.docs)
        assert "Pre-image text" in all_text
        assert "Post-image text" in all_text
        assert self.OCR_PREFIX in all_text

    def test_pptx_with_image_no_ocr_skips_images(self, tmp_path, mock_embedding):
        """Without OCR, images are skipped and only text is extracted."""
        from pptx import Presentation
        from pptx.util import Inches
        pptx_path = tmp_path / "img_skip.pptx"
        prs = Presentation()
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        txBox = slide.shapes.add_textbox(0, 0, 1000000, 1000000)
        txBox.text_frame.text = "Only text."
        img_bytes = _make_valid_png_bytes()
        slide.shapes.add_picture(io.BytesIO(img_bytes), Inches(1), Inches(2), Inches(2), Inches(2))
        prs.save(str(pptx_path))

        from bishon_kernel.core.local_file import LocalFile
        lf = LocalFile("user1", "KB001", str(pptx_path), "f3", "img_skip.pptx", mock_embedding)
        lf.split_file_to_docs(ocr_engine=None)
        all_text = " ".join(d.page_content for d in lf.docs)
        assert "Only text" in all_text
        assert self.OCR_PREFIX not in all_text


class TestDocxImageOcr:
    """Word image OCR: with ocr_engine, embedded image text is extracted."""

    OCR_PREFIX = "[DOCX_OCR]"

    def _mock_ocr(self, img_data):
        return [f"{self.OCR_PREFIX} ocr result"]

    def test_docx_without_ocr_still_extracts_text(self, tmp_path, mock_embedding):
        """Without OCR, Docx2txtLoader still extracts text normally."""
        from docx import Document as DocxDoc
        docx_path = tmp_path / "text_only.docx"
        doc = DocxDoc()
        doc.add_paragraph("Hello from docx without images.")
        doc.save(str(docx_path))

        from bishon_kernel.core.local_file import LocalFile
        lf = LocalFile("user1", "KB001", str(docx_path), "f1", "text_only.docx", mock_embedding)
        lf.split_file_to_docs(ocr_engine=None)
        all_text = " ".join(d.page_content for d in lf.docs)
        assert "Hello from docx without images" in all_text

    def test_docx_with_image_ocr(self, tmp_path, mock_embedding):
        """With OCR, embedded image text is extracted."""
        from docx import Document as DocxDoc
        from docx.shared import Inches
        docx_path = tmp_path / "with_image.docx"
        doc = DocxDoc()
        doc.add_paragraph("Text before image.")
        img_bytes = _make_valid_png_bytes()
        doc.add_picture(io.BytesIO(img_bytes), width=Inches(1))
        doc.add_paragraph("Text after image.")
        doc.save(str(docx_path))

        from bishon_kernel.core.local_file import LocalFile
        lf = LocalFile("user1", "KB001", str(docx_path), "f2", "with_image.docx", mock_embedding)
        lf.split_file_to_docs(ocr_engine=self._mock_ocr)
        all_text = " ".join(d.page_content for d in lf.docs)
        assert "Text before image" in all_text
        assert "Text after image" in all_text
        assert self.OCR_PREFIX in all_text

    def test_docx_with_image_no_ocr_skips_images(self, tmp_path, mock_embedding):
        """Without OCR, images are skipped and only text is extracted."""
        from docx import Document as DocxDoc
        from docx.shared import Inches
        docx_path = tmp_path / "img_skip.docx"
        doc = DocxDoc()
        doc.add_paragraph("Plain text only.")
        img_bytes = _make_valid_png_bytes()
        doc.add_picture(io.BytesIO(img_bytes), width=Inches(1))
        doc.save(str(docx_path))

        from bishon_kernel.core.local_file import LocalFile
        lf = LocalFile("user1", "KB001", str(docx_path), "f3", "img_skip.docx", mock_embedding)
        lf.split_file_to_docs(ocr_engine=None)
        all_text = " ".join(d.page_content for d in lf.docs)
        assert "Plain text only" in all_text
        assert self.OCR_PREFIX not in all_text
